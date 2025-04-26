# -*- coding: utf-8 -*-
import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt # Inches/Pt might be needed for positioning/sizing
import zipfile
from pathlib import Path
from PIL import Image, ImageTk
import tempfile
import traceback
import time # For unique temp names if needed

class AdvancedSlidesUtilitiesApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Advanced Slides Utilities")
        self.root.geometry("700x650") # Increased height slightly for status area
        self.root.minsize(650, 600)

        # --- Style Configuration ---
        self.style = ttk.Style()
        try:
            if os.name == 'nt': self.style.theme_use('vista')
            elif os.uname().sysname == 'Darwin': self.style.theme_use('aqua')
            else: self.style.theme_use('clam')
        except Exception: # Broad exception for safety on unusual systems
            self.style.theme_use('default')

        # Style for the Process Button (Red background, Blue text)
        self.style.configure('Action.TButton', font=('Arial', 12, 'bold'), foreground='blue', background='red', padding=10)
        self.style.map('Action.TButton',
            foreground=[('pressed', 'blue'), ('active', 'blue'), ('disabled', 'grey')],
            background=[('pressed', '!disabled', '#CC0000'), ('active', '#E50000'), ('disabled', '#F0F0F0')]
        )
        # Style for standard buttons
        self.style.configure('TButton', padding=5)
        self.style.configure('TEntry', padding=5)
        self.style.configure('secondary.TLabel', foreground='grey') # Style for hint text

        # --- Variables ---
        self.input_path_display = tk.StringVar() # Display single zip or count of pptx
        self.output_path = tk.StringVar()
        self.bg_image_path = tk.StringVar()
        self.use_bg_image = tk.BooleanVar(value=False)
        self.input_type = None # 'pptx', 'zip', None
        self.input_files = [] # List of pptx paths OR single zip path

        # --- Main Frame ---
        main_frame = ttk.Frame(root, padding="15 15 15 15")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # --- Header ---
        ttk.Label(main_frame, text="Advanced Slides Utilities", font=("Arial", 20, "bold"), anchor=tk.CENTER).pack(pady=(0, 20), fill=tk.X)

        # --- Input Frame ---
        input_frame = ttk.LabelFrame(main_frame, text=" 1. Select Input (PPTX file(s) or one ZIP) ", padding="10")
        input_frame.pack(fill=tk.X, pady=10)

        self.input_entry = ttk.Entry(input_frame, textvariable=self.input_path_display, width=60, state="readonly")
        self.input_entry.grid(row=0, column=0, padx=(0, 10), pady=5, sticky="ew")
        self.browse_input_btn = ttk.Button(input_frame, text="Browse...", command=self.browse_input, style='TButton')
        self.browse_input_btn.grid(row=0, column=1, padx=5, pady=5)
        input_frame.columnconfigure(0, weight=1)

        # --- Output Frame ---
        output_frame = ttk.LabelFrame(main_frame, text=" 2. Choose Output Destination ", padding="10")
        output_frame.pack(fill=tk.X, pady=10)

        self.output_label = ttk.Label(output_frame, text="Output Path:")
        self.output_label.grid(row=0, column=0, columnspan=2, padx=5, pady=(0,5), sticky=tk.W)

        self.output_entry = ttk.Entry(output_frame, textvariable=self.output_path, width=60)
        self.output_entry.grid(row=1, column=0, padx=(0, 10), pady=5, sticky="ew")
        self.browse_output_btn = ttk.Button(output_frame, text="Browse...", command=self.browse_output, style='TButton')
        self.browse_output_btn.grid(row=1, column=1, padx=5, pady=5)
        output_frame.columnconfigure(0, weight=1)

        # --- Background Options Frame ---
        bg_frame = ttk.LabelFrame(main_frame, text=" 3. Background Options (Optional) ", padding="10")
        bg_frame.pack(fill=tk.X, pady=10)

        self.bg_check = ttk.Checkbutton(bg_frame, text="Apply Custom Background Image", variable=self.use_bg_image, command=self.toggle_bg_image)
        self.bg_check.grid(row=0, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)

        self.bg_image_entry = ttk.Entry(bg_frame, textvariable=self.bg_image_path, width=55, state="disabled")
        self.bg_image_entry.grid(row=1, column=0, padx=(5, 10), pady=5, sticky="ew")
        self.bg_image_button = ttk.Button(bg_frame, text="Select Image...", command=self.browse_bg_image, state="disabled", style='TButton')
        self.bg_image_button.grid(row=1, column=1, padx=5, pady=5)
        ttk.Label(bg_frame, text="(If unchecked, master backgrounds set to white. Overwrites existing backgrounds.)", font=("Arial", 9), style='secondary.TLabel').grid(row=2, column=0, columnspan=2, padx=5, pady=(5,0), sticky=tk.W)
        bg_frame.columnconfigure(0, weight=1)

        # --- Separator and Action Button ---
        ttk.Separator(main_frame, orient=tk.HORIZONTAL).pack(fill=tk.X, pady=20)

        # --- Process Button ---
        self.process_button = ttk.Button(main_frame, text="Process Slides", command=self.process_pptx, style='Action.TButton', width=20)
        self.process_button.pack(pady=10)

        # --- Status Label ---
        self.status_label = ttk.Label(main_frame, text="Ready", wraplength=650, anchor=tk.CENTER, justify=tk.CENTER, font=("Arial", 10))
        self.status_label.pack(pady=(15, 0), fill=tk.X, expand=True)


    # --- Browse and Selection Logic ---

    def browse_input(self):
        """Allows selecting multiple PPTX files OR a single ZIP file."""
        files = filedialog.askopenfilenames(
            title="Select Input PPTX File(s) OR a Single ZIP Archive",
            # Added ZIP to filetypes
            filetypes=[("PowerPoint/ZIP", "*.pptx *.zip"),
                       ("PowerPoint Files", "*.pptx"),
                       ("ZIP Archives", "*.zip"),
                       ("All Files", "*.*")]
        )
        if not files:
            return # User cancelled

        # Check selection: only pptx OR only one zip
        pptx_files = [f for f in files if f.lower().endswith('.pptx')]
        zip_files = [f for f in files if f.lower().endswith('.zip')]

        if zip_files and pptx_files:
            messagebox.showerror("Invalid Selection", "Please select EITHER PowerPoint (.pptx) files OR a single ZIP (.zip) file, not both.")
            return
        if len(zip_files) > 1:
            messagebox.showerror("Invalid Selection", "Please select only ONE ZIP file at a time.")
            return

        # Process valid selection
        first_file_dir = os.path.dirname(files[0])
        self.output_path.set(first_file_dir) # Suggest output dir

        if pptx_files:
            self.input_type = 'pptx'
            self.input_files = pptx_files
            if len(pptx_files) == 1:
                self.input_path_display.set(os.path.basename(pptx_files[0]))
                self.output_label.config(text="Output File Path:")
                base, ext = os.path.splitext(os.path.basename(pptx_files[0]))
                suggested_output = os.path.join(first_file_dir, f"{base}_processed{ext}")
                self.output_path.set(suggested_output)
            else:
                self.input_path_display.set(f"{len(pptx_files)} PPTX files selected")
                self.output_label.config(text="Output Directory (for ZIP):")
                self.output_path.set(first_file_dir) # Suggest dir for multi-pptx output zip

        elif zip_files:
            self.input_type = 'zip'
            self.input_files = zip_files # Store the single zip path
            self.input_path_display.set(os.path.basename(zip_files[0]))
            self.output_label.config(text="Output Directory (for processed ZIP):")
            self.output_path.set(first_file_dir) # Suggest dir for zip output

        else: # Should not happen if files were selected, but defensively:
            self.input_type = None
            self.input_files = []
            self.input_path_display.set("")
            messagebox.showwarning("No Files Found", "No valid .pptx or .zip files were selected.")


    def browse_output(self):
        if not self.input_files:
             messagebox.showwarning("Select Input First", "Please select input file(s) or ZIP before choosing the output.")
             return

        current_output_suggestion = self.output_path.get()
        initial_dir = os.path.dirname(current_output_suggestion) if current_output_suggestion and os.path.isdir(os.path.dirname(current_output_suggestion)) else (os.path.dirname(self.input_files[0]) if self.input_files else "/")

        if self.input_type == 'pptx' and len(self.input_files) > 1: # Multiple PPTX -> output dir for zip
            directory = filedialog.askdirectory(
                title="Select Output Directory for saving processed PPTX files (as ZIP)",
                initialdir=initial_dir
                )
            if directory:
                self.output_path.set(directory)
        elif self.input_type == 'zip': # ZIP input -> output dir for new zip
             directory = filedialog.askdirectory(
                title="Select Output Directory for saving the new ZIP archive",
                initialdir=initial_dir
                )
             if directory:
                self.output_path.set(directory)
        elif self.input_type == 'pptx' and len(self.input_files) == 1: # Single PPTX -> output file
            initial_file = os.path.basename(current_output_suggestion) if current_output_suggestion and not os.path.isdir(current_output_suggestion) else ""
            file_path = filedialog.asksaveasfilename(
                title="Save Processed PPTX As",
                defaultextension=".pptx",
                filetypes=[("PowerPoint Files", "*.pptx")],
                initialdir=initial_dir,
                initialfile=initial_file or f"processed_{os.path.basename(self.input_files[0])}"
            )
            if file_path:
                self.output_path.set(file_path)


    def browse_bg_image(self):
        # (No changes needed here from previous version)
        initial_dir = os.path.dirname(self.bg_image_path.get()) if self.bg_image_path.get() else "/"
        file_path = filedialog.askopenfilename(
            title="Select Background Image",
            filetypes=[("Image Files", "*.png *.jpg *.jpeg"), ("All Files", "*.*")],
            initialdir=initial_dir
        )
        if file_path:
            try:
                with Image.open(file_path) as img:
                    img_format = img.format
                    if img_format not in ["PNG", "JPEG"]:
                        messagebox.showerror("Invalid Format", f"Unsupported image format: {img_format}. Please select a PNG or JPEG.")
                        return
                    if os.path.getsize(file_path) > 15 * 1024 * 1024: # 15 MB limit warning
                        messagebox.showwarning("Large File", "Warning: Image file is large (>15 MB). This may significantly increase presentation size and processing time.")
                self.bg_image_path.set(file_path)
            except FileNotFoundError:
                 messagebox.showerror("Error", f"File not found: {file_path}")
            except Exception as e:
                messagebox.showerror("Image Error", f"Could not open or validate image file:\n{str(e)}")
                self.bg_image_path.set("")


    def toggle_bg_image(self):
        # (No changes needed here from previous version)
        state = "normal" if self.use_bg_image.get() else "disabled"
        # Make sure button isn't enabled if main process is running
        if self.process_button['state'] == 'disabled':
             state = 'disabled'
        self.bg_image_entry.config(state=state)
        self.bg_image_button.config(state=state)
        if not self.use_bg_image.get():
             self.bg_image_path.set("")


    # --- Core Processing Logic ---

    def set_background(self, prs, slide_master, bg_image_path=None):
        """
        Sets the background on a slide master.
        If bg_image_path is provided, adds it as a full-size picture at the back.
        Otherwise, sets a solid white background fill.
        Requires the presentation object (prs) for dimensions.
        """
        background = slide_master.background
        fill = background.fill

        # Clear existing shapes that might be acting as backgrounds (optional, can be risky)
        # for shape in reversed(slide_master.shapes):
        #     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.width == prs.slide_width and shape.height == prs.slide_height:
        #          try:
        #             sp = shape._element
        #             sp.getparent().remove(sp)
        #             print(f"DEBUG: Removed potential old background picture shape from master.")
        #          except Exception: pass

        if bg_image_path and os.path.exists(bg_image_path):
            try:
                # Set background fill to solid white first to clear any previous fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)

                # Add the picture as a shape covering the whole slide
                left, top = 0, 0
                width, height = prs.slide_width, prs.slide_height
                pic = slide_master.shapes.add_picture(bg_image_path, left, top, width=width, height=height)

                # Send the picture to the back
                slide_master.shapes._spTree.remove(pic._element)
                slide_master.shapes._spTree.insert(2, pic._element) # Index 2 typically places it behind content

            except Exception as e:
                 print(f"ERROR setting picture background for master: {e}\n{traceback.format_exc()}")
                 self.status_label.config(text=f"Warning: Failed to set picture background. Check image.")
                 # Fallback to white
                 fill.solid()
                 fill.fore_color.rgb = RGBColor(255, 255, 255)
        else:
            # Set solid white background fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)


    def optimize_pptx(self, input_path, output_path, bg_image_path_to_use):
        """ Processes a single PPTX file. Returns (orig_mb, new_mb) or (None, None). """
        temp_prs_path = None
        try:
            self.status_label.config(text=f"Processing {os.path.basename(input_path)}...")
            self.root.update_idletasks()

            prs = Presentation(input_path)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # --- Background Application ---
            # 1. Apply to Slide Masters
            for master in prs.slide_masters:
                self.set_background(prs, master, bg_image_path_to_use) # Pass prs for dimensions

            # 2. Clear Individual Slide Backgrounds (Force Inheritance)
            for slide in prs.slides:
                slide.background.fill.background() # Inherit from master

            # 3. Set Slide Layout Backgrounds to White (Safety Net)
            for layout in prs.slide_layouts:
                 layout.background.fill.solid()
                 layout.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # --- Save and Re-compress ---
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                 temp_prs_path = temp_file.name
            prs.save(temp_prs_path)

            with zipfile.ZipFile(temp_prs_path, 'r') as zin:
                with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
                    for item in zin.infolist():
                        buffer = zin.read(item.filename)
                        zout.writestr(item, buffer)

            orig_size_bytes = os.path.getsize(input_path)
            new_size_bytes = os.path.getsize(output_path)
            return orig_size_bytes / (1024 * 1024), new_size_bytes / (1024 * 1024)

        except Exception as e:
            error_msg = f"Error processing {os.path.basename(input_path)}: {str(e)}"
            print(f"--- ERROR --- \n{error_msg}\n{traceback.format_exc()}\n-------------")
            self.status_label.config(text=error_msg)
            if os.path.exists(output_path):
                 try: os.remove(output_path)
                 except Exception as del_err: print(f"Warn: Could not remove failed output {output_path}: {del_err}")
            return None, None
        finally:
            if temp_prs_path and os.path.exists(temp_prs_path):
                try: os.remove(temp_prs_path)
                except OSError as e: print(f"Warn: Could not delete temp file {temp_prs_path}: {e}")


    def process_pptx(self):
        """ Main processing controller triggered by the button. Handles single/multi PPTX and ZIP input. """
        # --- Validations ---
        if not self.input_files or self.input_type is None:
            messagebox.showerror("Input Missing", "Please select input file(s) or a ZIP archive.")
            return
        output_dest = self.output_path.get()
        if not output_dest:
            messagebox.showerror("Output Missing", "Please select an output destination.")
            return

        bg_image_path_selected = self.bg_image_path.get()
        apply_custom_bg = self.use_bg_image.get()
        final_bg_image_path = None

        if apply_custom_bg:
            if not bg_image_path_selected:
                messagebox.showerror("BG Image Missing", "Custom background checked, but no image selected.")
                return
            if not os.path.exists(bg_image_path_selected):
                 messagebox.showerror("BG Image Invalid", f"BG image not found:\n{bg_image_path_selected}")
                 return
            final_bg_image_path = bg_image_path_selected

        # --- Disable UI ---
        self.process_button.config(state="disabled")
        self.browse_input_btn.config(state="disabled")
        self.browse_output_btn.config(state="disabled")
        self.bg_check.config(state="disabled")
        self.bg_image_button.config(state="disabled")
        self.status_label.config(text="Starting processing...")
        self.root.update()

        # --- Processing Logic ---
        start_time = time.time()
        try:
            if self.input_type == 'pptx':
                self.process_pptx_files(self.input_files, output_dest, final_bg_image_path)
            elif self.input_type == 'zip':
                self.process_zip_file(self.input_files[0], output_dest, final_bg_image_path)

        except Exception as general_e:
            error_msg = f"An unexpected error occurred: {str(general_e)}"
            print(f"--- UNEXPECTED ERROR ---\n{error_msg}\n{traceback.format_exc()}\n------------------------")
            self.status_label.config(text=error_msg)
            messagebox.showerror("Unexpected Error", f"An error occurred:\n{general_e}")

        finally:
            # --- Re-enable UI ---
            end_time = time.time()
            print(f"Total processing time: {end_time - start_time:.2f} seconds")
            self.process_button.config(state="normal")
            self.browse_input_btn.config(state="normal")
            self.browse_output_btn.config(state="normal")
            self.bg_check.config(state="normal")
            self.toggle_bg_image() # Correctly sets bg button state based on check
            self.root.update()


    def process_pptx_files(self, pptx_files_list, output_destination, bg_image_path):
        """ Handles processing for one or more PPTX files. """
        num_input_files = len(pptx_files_list)

        if num_input_files == 1:
            # --- Single PPTX File ---
            input_path = pptx_files_list[0]
            output_path = output_destination

            if os.path.isdir(output_path): # User selected dir instead of file
                 output_filename = f"processed_{os.path.basename(input_path)}"
                 output_path = os.path.join(output_path, output_filename)
                 self.output_path.set(output_path) # Update UI

            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            orig_size, new_size = self.optimize_pptx(input_path, output_path, bg_image_path)

            if orig_size is not None:
                reduction = orig_size - new_size
                percent = (reduction / orig_size * 100) if orig_size > 0 else 0
                final_msg = (f"Success! File saved to:\n{os.path.basename(output_path)}\n"
                             f"Original: {orig_size:.2f} MB, Processed: {new_size:.2f} MB ({percent:+.1f}%)")
                self.status_label.config(text=final_msg)
                messagebox.showinfo("Processing Complete", final_msg)
            else:
                fail_msg = f"Failed to process:\n{os.path.basename(input_path)}"
                self.status_label.config(text=f"FAILED: {os.path.basename(input_path)}. Check console.")
                messagebox.showerror("Processing Failed", f"{fail_msg}\nSee status bar or console for error details.")

        else:
            # --- Multiple PPTX Files (Output to ZIP) ---
            output_dir = output_destination
            if not os.path.isdir(output_dir):
                messagebox.showerror("Output Directory Needed", "For multiple PPTX files, the output must be a directory to save the ZIP file.")
                # No 'return' here, error handled in finally block re-enabling UI
                self.status_label.config(text="Error: Output destination must be a directory for multiple files.")
                return # Explicitly return here after setting status

            zip_filename = "processed_slides.zip"
            zip_path = os.path.join(output_dir, zip_filename)
            processed_count = 0
            failed_count = 0
            total_orig_mb = 0
            total_new_mb = 0

            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    with zipfile.ZipFile(zip_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
                        for i, input_path in enumerate(pptx_files_list):
                            # Use a unique temp name inside the temp dir
                            temp_output_filename = f"{i}_{os.path.basename(input_path)}"
                            temp_output_path = os.path.join(temp_dir, temp_output_filename)

                            orig_size, new_size = self.optimize_pptx(input_path, temp_output_path, bg_image_path)

                            if orig_size is not None and os.path.exists(temp_output_path):
                                base, ext = os.path.splitext(os.path.basename(input_path))
                                archive_name = f"{base}_processed{ext}"
                                zf.write(temp_output_path, arcname=archive_name)
                                processed_count += 1
                                total_orig_mb += orig_size
                                total_new_mb += new_size
                            else:
                                failed_count += 1

                    # --- After processing all files ---
                    if processed_count > 0:
                        final_zip_size_mb = os.path.getsize(zip_path) / (1024 * 1024)
                        status_msg = (
                            f"Processed {processed_count}/{num_input_files} files. Saved to: {zip_filename}\n"
                            f"Total Original: {total_orig_mb:.2f} MB, Total Processed (est.): {total_new_mb:.2f} MB.\n"
                            f"Final ZIP size: {final_zip_size_mb:.2f} MB."
                        )
                        if failed_count > 0:
                            status_msg += f"\nWarning: {failed_count} file(s) failed."
                        self.status_label.config(text=status_msg)
                        messagebox.showinfo("Processing Complete", status_msg)
                    elif failed_count > 0:
                        fail_msg = f"Processing failed. {failed_count} file(s) could not be processed."
                        self.status_label.config(text=fail_msg)
                        if os.path.exists(zip_path): os.remove(zip_path)
                        messagebox.showerror("Processing Failed", f"{fail_msg}\nNo ZIP file created.")
                    else:
                         self.status_label.config(text="No files were processed successfully.")
                         if os.path.exists(zip_path): os.remove(zip_path)

                except Exception as zip_e:
                     error_msg = f"Fatal error creating ZIP: {str(zip_e)}"
                     print(f"{error_msg}\n{traceback.format_exc()}")
                     self.status_label.config(text=error_msg)
                     messagebox.showerror("ZIP Error", f"Could not create ZIP:\n{zip_e}")
                     if os.path.exists(zip_path):
                         try: os.remove(zip_path)
                         except Exception as del_err: print(f"Warn: Could not remove partial zip {zip_path}: {del_err}")


    def process_zip_file(self, zip_input_path, output_dir, bg_image_path):
        """ Handles processing when the input is a single ZIP archive. """
        if not os.path.isdir(output_dir):
            messagebox.showerror("Output Directory Required", "For ZIP input, the output destination must be an existing directory.")
            self.status_label.config(text="Error: Output destination must be a directory for ZIP input.")
            return

        extracted_files = []
        processed_count = 0
        failed_count = 0
        total_orig_mb = 0
        total_new_mb = 0

        output_zip_filename = "processed_slides_from_zip.zip"
        output_zip_path = os.path.join(output_dir, output_zip_filename)

        # Use two nested temp dirs: one for extraction, one for processing results before zipping
        with tempfile.TemporaryDirectory() as extract_temp_dir, tempfile.TemporaryDirectory() as process_temp_dir:
            try:
                # --- 1. Extract PPTX from Input ZIP ---
                self.status_label.config(text=f"Extracting PPTX files from {os.path.basename(zip_input_path)}...")
                self.root.update_idletasks()
                num_extracted = 0
                with zipfile.ZipFile(zip_input_path, 'r') as zin:
                    for item in zin.infolist():
                        if not item.is_dir() and item.filename.lower().endswith('.pptx'):
                            # Ensure we handle potential subdirectories in the zip
                            target_path = os.path.join(extract_temp_dir, item.filename)
                            os.makedirs(os.path.dirname(target_path), exist_ok=True)
                            try:
                                zin.extract(item, path=extract_temp_dir)
                                extracted_files.append(target_path) # Store full path to extracted file
                                num_extracted += 1
                            except Exception as extract_err:
                                print(f"Warning: Failed to extract {item.filename} from zip: {extract_err}")
                self.status_label.config(text=f"Extracted {num_extracted} PPTX files. Starting processing...")
                self.root.update_idletasks()

                if not extracted_files:
                     self.status_label.config(text="No .pptx files found inside the selected ZIP archive.")
                     messagebox.showwarning("No Files Found", "The selected ZIP archive does not contain any .pptx files.")
                     return

                # --- 2. Process Extracted Files (like multi-pptx case) ---
                with zipfile.ZipFile(output_zip_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
                    for i, input_path in enumerate(extracted_files):
                         # Define unique temp path for the *processed* version in the *second* temp dir
                        temp_output_filename = f"{i}_{os.path.basename(input_path)}"
                        temp_output_path = os.path.join(process_temp_dir, temp_output_filename)

                        orig_size, new_size = self.optimize_pptx(input_path, temp_output_path, bg_image_path)

                        if orig_size is not None and os.path.exists(temp_output_path):
                            # Use the original filename from *within the zip* as the base for the archive name
                            original_arcname_base = Path(input_path).relative_to(extract_temp_dir) # Get path relative to temp dir
                            base, ext = os.path.splitext(original_arcname_base)
                            archive_name = f"{base}_processed{ext}"
                            zf.write(temp_output_path, arcname=str(archive_name)) # arcname needs to be string
                            processed_count += 1
                            total_orig_mb += orig_size
                            total_new_mb += new_size
                        else:
                            failed_count += 1

                # --- 3. Report Results ---
                if processed_count > 0:
                    final_zip_size_mb = os.path.getsize(output_zip_path) / (1024 * 1024)
                    status_msg = (
                        f"Processed {processed_count}/{len(extracted_files)} PPTX from ZIP. Saved to: {output_zip_filename}\n"
                        f"Total Original (est.): {total_orig_mb:.2f} MB, Total Processed (est.): {total_new_mb:.2f} MB.\n"
                        f"Final ZIP size: {final_zip_size_mb:.2f} MB."
                    )
                    if failed_count > 0:
                        status_msg += f"\nWarning: {failed_count} extracted file(s) failed processing."
                    self.status_label.config(text=status_msg)
                    messagebox.showinfo("Processing Complete", status_msg)
                elif failed_count > 0:
                    fail_msg = f"Processing failed. {failed_count} extracted file(s) could not be processed."
                    self.status_label.config(text=fail_msg)
                    if os.path.exists(output_zip_path): os.remove(output_zip_path)
                    messagebox.showerror("Processing Failed", f"{fail_msg}\nNo output ZIP file created.")
                else: # Should only happen if extraction worked but all processing failed
                     self.status_label.config(text="All extracted PPTX files failed during processing.")
                     if os.path.exists(output_zip_path): os.remove(output_zip_path)

            except Exception as zip_proc_e:
                 error_msg = f"Error processing ZIP archive: {str(zip_proc_e)}"
                 print(f"{error_msg}\n{traceback.format_exc()}")
                 self.status_label.config(text=error_msg)
                 messagebox.showerror("ZIP Processing Error", f"Could not process the ZIP file:\n{zip_proc_e}")
                 if os.path.exists(output_zip_path):
                     try: os.remove(output_zip_path)
                     except Exception as del_err: print(f"Warn: Could not remove partial output zip {output_zip_path}: {del_err}")


# --- Main Execution ---
if __name__ == "__main__":
    root = tk.Tk()
    app = AdvancedSlidesUtilitiesApp(root)
    root.mainloop()
